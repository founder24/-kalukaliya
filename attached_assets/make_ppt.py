"""
Generate a 2-slide technical block diagram PPT for MSME IdeaHackathon 6.0.
Slide 1 — System Architecture Block Diagram (functional modules + data flows)
Slide 2 — Technical Processing Brief (pipeline stages + parameters)
No identifying information anywhere.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0E, 0x1A)   # near-black navy
PANEL     = RGBColor(0x11, 0x18, 0x27)   # card bg
BORDER    = RGBColor(0x1E, 0x2A, 0x4A)   # card border (not directly usable in pptx fill)
INDIGO    = RGBColor(0x4F, 0x46, 0xE5)   # board level
BLUE      = RGBColor(0x3B, 0x82, 0xF6)   # class level
CYAN      = RGBColor(0x06, 0xB6, 0xD4)   # stream/language
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # subject/RAG
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # chapter/LLM
PINK      = RGBColor(0xEC, 0x48, 0x99)   # topic/output
PURPLE    = RGBColor(0x8B, 0x5C, 0xF6)   # memory
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LTGREY    = RGBColor(0xC8, 0xD0, 0xE8)
MIDGREY   = RGBColor(0x6B, 0x7D, 0xB3)
DARKGREY  = RGBColor(0x1E, 0x29, 0x47)

# ── Slide dimensions: 16 × 9 widescreen ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width_pt=1.0):
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width_pt=1.2, radius=0.08):
    """MSO_SHAPE 5 = ROUNDED_RECTANGLE"""
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    # pptx adjusts[0] controls corner radius (0=square, 50000=half-round)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, l, t, w, h, text, font_size=9, bold=False,
                 color=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_arrow_h(slide, lx, ty, length, color=MIDGREY, width_pt=1.5):
    """Horizontal arrow left→right"""
    from pptx.util import Inches
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(lx), Inches(ty),
        Inches(lx + length), Inches(ty)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)


def add_arrow_v(slide, lx, ty, length, color=MIDGREY, width_pt=1.5):
    """Vertical arrow top→bottom"""
    connector = slide.shapes.add_connector(
        1,
        Inches(lx), Inches(ty),
        Inches(lx), Inches(ty + length)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)


def label_box(slide, l, t, w, h, title, subtitle, fill, border, title_size=8, sub_size=6.5):
    """Filled rounded box with a bold title line and a smaller subtitle line."""
    add_rounded_rect(slide, l, t, w, h, fill, border, line_width_pt=1.4)
    # title
    add_text_box(slide, l, t + 0.02, w, h * 0.5,
                 title, font_size=title_size, bold=True, color=WHITE)
    # subtitle
    if subtitle:
        add_text_box(slide, l, t + h * 0.45, w, h * 0.5,
                     subtitle, font_size=sub_size, bold=False, color=LTGREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — SYSTEM ARCHITECTURE BLOCK DIAGRAM
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)

# Background
add_rect(s1, 0, 0, 13.33, 7.5, BG)

# ── Header bar ────────────────────────────────────────────────────────────────
add_rect(s1, 0, 0, 13.33, 0.52, DARKGREY)
add_text_box(s1, 0.2, 0.04, 13, 0.42,
             "BILINGUAL AI EDUCATIONAL PLATFORM  ·  SYSTEM ARCHITECTURE BLOCK DIAGRAM",
             font_size=11, bold=True, color=LTGREY)

# ── Section label: top-left (Hierarchy) ───────────────────────────────────────
add_text_box(s1, 0.18, 0.58, 3.4, 0.25,
             "CONTENT HIERARCHY  (Library Module)", font_size=7, bold=True, color=INDIGO)

# Hierarchy boxes — vertical stack left side
hier = [
    ("BOARD",    "SEBA · AHSEC · CBSE",          INDIGO, 0.72, RGBColor(0x31,0x2E,0x81)),
    ("CLASS",    "Grade 9 / 10 / 11 / 12",       BLUE,   1.22, RGBColor(0x1E,0x3A,0x5F)),
    ("STREAM",   "Science · Arts · Commerce",     CYAN,   1.72, RGBColor(0x0C,0x2D,0x35)),
    ("SUBJECT",  "Physics · Chem · Math · Bio…",  GREEN,  2.22, RGBColor(0x0A,0x20,0x18)),
    ("CHAPTER",  "Chapter 1 … Chapter N",         AMBER,  2.72, RGBColor(0x20,0x15,0x00)),
    ("TOPIC",    "Topic Units (Granular)",         PINK,   3.22, RGBColor(0x24,0x0A,0x18)),
]

for (lv, sub, col, ty, dark) in hier:
    add_rounded_rect(s1, 0.18, ty, 1.65, 0.40, dark, col, 1.3)
    add_text_box(s1, 0.18, ty + 0.02, 1.65, 0.20, lv, font_size=8, bold=True, color=WHITE)
    add_text_box(s1, 0.18, ty + 0.20, 1.65, 0.18, sub, font_size=6, bold=False, color=LTGREY)

# vertical connector between hierarchy boxes
for i in range(len(hier) - 1):
    ty = hier[i][3] + 0.40
    add_arrow_v(s1, 1.0, ty, 0.10, color=hier[i][2], width_pt=1.5)

# Content types capsule
add_rounded_rect(s1, 0.18, 3.72, 1.65, 0.52, RGBColor(0x12,0x0B,0x2A), PURPLE, 1.2)
add_text_box(s1, 0.18, 3.74, 1.65, 0.22, "CONTENT TYPES", font_size=7, bold=True, color=PURPLE)
add_text_box(s1, 0.18, 3.94, 1.65, 0.28,
             "Notes (EN/AS)  ·  Q&A (EN/AS)\nPYQ Papers  ·  Voice Output",
             font_size=6, bold=False, color=LTGREY)

# arrow from hierarchy → pipeline
add_arrow_h(s1, 1.85, 2.50, 0.42, color=CYAN, width_pt=1.8)


# ── CENTRE: Processing Pipeline (horizontal flow) ─────────────────────────────
add_text_box(s1, 2.35, 0.58, 10.8, 0.25,
             "QUERY PROCESSING PIPELINE", font_size=7, bold=True, color=CYAN)

pipe_y  = 0.88   # top of first pipeline row
pipe_h  = 0.72
gap     = 0.08
bw      = 1.42   # box width
# 7 pipeline stages
stages = [
    ("INPUT\nINTERFACE",    "Text / Voice\nEN or AS",           INDIGO,  RGBColor(0x1E,0x1B,0x4B)),
    ("LANGUAGE\nDETECTOR",  "Script analysis\n+ LID model",     BLUE,    RGBColor(0x1E,0x3A,0x5F)),
    ("QUERY\nEMBEDDER",     "BGE-M3 1024-dim\nmultilingual",    CYAN,    RGBColor(0x0C,0x2D,0x35)),
    ("TOPIC\nMATCHER",      "Cosine similarity\n197 embeddings", GREEN,   RGBColor(0x0A,0x20,0x18)),
    ("RAG\nRETRIEVAL",      "Top-K chunks\nVector DB lookup",    AMBER,   RGBColor(0x20,0x15,0x00)),
    ("LLM\nGENERATOR",      "Gemini / Indic\nLLM (lang-aware)", PINK,    RGBColor(0x24,0x0A,0x18)),
    ("RESPONSE\nOUTPUT",    "Streamed text\n+ TTS voice",        PURPLE,  RGBColor(0x1A,0x08,0x40)),
]

start_x = 2.30
for i, (title, sub, col, dark) in enumerate(stages):
    lx = start_x + i * (bw + gap)
    add_rounded_rect(s1, lx, pipe_y, bw, pipe_h, dark, col, 1.4)
    add_text_box(s1, lx, pipe_y + 0.04, bw, 0.34,
                 title, font_size=8, bold=True, color=WHITE)
    add_text_box(s1, lx, pipe_y + 0.36, bw, 0.32,
                 sub, font_size=6.5, bold=False, color=LTGREY)
    # arrow between stages
    if i < len(stages) - 1:
        ax = lx + bw + 0.01
        ay = pipe_y + pipe_h / 2
        add_arrow_h(s1, ax, ay, gap - 0.01, color=col, width_pt=1.6)

# ── DATA PARAMETERS ROW (under pipeline) ─────────────────────────────────────
params = [
    ("Input params",     "lang={en|as}\nmodality={text|voice}"),
    ("Detection output", "lang_code\nconfidence_score"),
    ("Embedding output", "vector[1024]\nquery_repr"),
    ("Matcher output",   "topic_id\nsimilarity_score"),
    ("Retrieval params", "top_k=5\nchunks + metadata"),
    ("Generator params", "prompt+context\nmax_tokens=800"),
    ("Output params",    "stream=SSE\nvoice={sarvam-tts}"),
]

param_y = pipe_y + pipe_h + 0.06
for i, (label, val) in enumerate(params):
    lx = start_x + i * (bw + gap)
    add_rounded_rect(s1, lx, param_y, bw, 0.50, RGBColor(0x0C,0x11,0x1E), BORDER, 0.8)
    add_text_box(s1, lx, param_y + 0.02, bw, 0.18,
                 label, font_size=6, bold=True, color=MIDGREY)
    add_text_box(s1, lx, param_y + 0.20, bw, 0.26,
                 val, font_size=6.5, bold=False, color=LTGREY)


# ── MEMORY MODULE ─────────────────────────────────────────────────────────────
mem_x, mem_y, mem_w, mem_h = 2.30, 2.30, 5.5, 0.60
add_rounded_rect(s1, mem_x, mem_y, mem_w, mem_h, RGBColor(0x12,0x0B,0x2A), PURPLE, 1.3)
add_text_box(s1, mem_x, mem_y + 0.03, mem_w, 0.26,
             "PERSISTENT STUDENT MEMORY MODULE", font_size=8, bold=True, color=PURPLE)
add_text_box(s1, mem_x, mem_y + 0.30, mem_w, 0.26,
             "Stores: topics_studied[], weak_areas[], doubt_history[], session_count  ·  Read on every query  ·  Written after every session",
             font_size=6.5, bold=False, color=LTGREY)

# arrow from pipeline down to memory
add_arrow_v(s1, 5.0, pipe_y + pipe_h, 0.30, color=PURPLE, width_pt=1.5)

# ── KNOWLEDGE BASE MODULE ─────────────────────────────────────────────────────
kb_x, kb_y, kb_w, kb_h = 8.10, 2.30, 5.05, 0.60
add_rounded_rect(s1, kb_x, kb_y, kb_w, kb_h, RGBColor(0x05,0x18,0x10), GREEN, 1.3)
add_text_box(s1, kb_x, kb_y + 0.03, kb_w, 0.26,
             "KNOWLEDGE BASE  (Vector Store + Document DB)", font_size=8, bold=True, color=GREEN)
add_text_box(s1, kb_x, kb_y + 0.30, kb_w, 0.26,
             "Indexed: chapter_notes, qa_pairs, pyq_text, formula_sheets  ·  Languages: EN + AS  ·  Embeddings: BGE-M3 (1024-dim)",
             font_size=6.5, bold=False, color=LTGREY)

# arrow from RAG retrieval up to KB
add_arrow_v(s1, 9.8, pipe_y + pipe_h, 0.30, color=GREEN, width_pt=1.5)


# ── ADMIN / STAFF CONTENT PIPELINE ───────────────────────────────────────────
add_text_box(s1, 2.30, 3.18, 10.85, 0.22,
             "CONTENT INGESTION PIPELINE (Staff Module → Knowledge Base)", font_size=7, bold=True, color=AMBER)

cp_y   = 3.42
cp_h   = 0.50
cp_bw  = 1.75
cp_gap = 0.12
cp_stages = [
    ("RAW CONTENT\nINGEST",    "PDF · Markdown\nManual entry"),
    ("PARSER /\nCHUNKER",     "Section split\nclean HTML→text"),
    ("BILINGUAL\nTRANSLATOR", "EN→AS pipeline\nSarvam translate"),
    ("EMBED +\nINDEX",        "BGE-M3 encode\nVectorize write"),
    ("PUBLISH\nGATEKEEPER",   "Staff QC review\nversion control"),
    ("LIVE IN\nKNOWLEDGE BASE","Immediately\navailable in RAG"),
]
cpx = 2.30
for i, (title, sub) in enumerate(cp_stages):
    lx = cpx + i * (cp_bw + cp_gap)
    add_rounded_rect(s1, lx, cp_y, cp_bw, cp_h, RGBColor(0x18,0x10,0x00), AMBER, 1.2)
    add_text_box(s1, lx, cp_y + 0.03, cp_bw, 0.22,
                 title, font_size=7.5, bold=True, color=WHITE)
    add_text_box(s1, lx, cp_y + 0.26, cp_bw, 0.22,
                 sub, font_size=6, bold=False, color=LTGREY)
    if i < len(cp_stages) - 1:
        ax = lx + cp_bw + 0.01
        ay = cp_y + cp_h / 2
        add_arrow_h(s1, ax, ay, cp_gap - 0.01, color=AMBER, width_pt=1.5)


# ── CURRICULUM RESOLVER BLOCK ──────────────────────────────────────────────────
add_text_box(s1, 2.30, 4.08, 10.85, 0.22,
             "CURRICULUM RESOLVER  —  Topic → Chapter → Subject → Stream → Class → Board  (automatic back-mapping on each query)",
             font_size=7.5, bold=True, color=CYAN)

resolver_levels = [
    ("TOPIC ID",   PINK),
    ("CHAPTER",    AMBER),
    ("SUBJECT",    GREEN),
    ("STREAM",     CYAN),
    ("CLASS",      BLUE),
    ("BOARD",      INDIGO),
    ("ANSWER\nDELIVERED", PURPLE),
]
rv_bw  = 1.42
rv_gap = 0.08
rv_y   = 4.35
rv_h   = 0.44
rvx    = 2.30
for i, (label, col) in enumerate(resolver_levels):
    lx = rvx + i * (rv_bw + rv_gap)
    dark_map = {PINK: RGBColor(0x24,0x0A,0x18), AMBER: RGBColor(0x20,0x15,0x00),
                GREEN: RGBColor(0x0A,0x20,0x18), CYAN: RGBColor(0x0C,0x2D,0x35),
                BLUE: RGBColor(0x1E,0x3A,0x5F), INDIGO: RGBColor(0x1E,0x1B,0x4B),
                PURPLE: RGBColor(0x1A,0x08,0x40)}
    add_rounded_rect(s1, lx, rv_y, rv_bw, rv_h, dark_map[col], col, 1.3)
    add_text_box(s1, lx, rv_y + 0.10, rv_bw, 0.26,
                 label, font_size=8, bold=True, color=WHITE)
    if i < len(resolver_levels) - 1:
        ax = lx + rv_bw + 0.01
        ay = rv_y + rv_h / 2
        add_arrow_h(s1, ax, ay, rv_gap - 0.01, color=col, width_pt=1.6)


# ── TECH STACK STRIP ─────────────────────────────────────────────────────────
stack_y = 4.94
add_text_box(s1, 0.18, stack_y, 2.0, 0.22,
             "TECHNOLOGY STACK", font_size=7, bold=True, color=MIDGREY)

stack_items = [
    ("Edge Runtime",    "Cloudflare Workers"),
    ("Vector DB",       "Vectorize (HNSW)"),
    ("Embed Model",     "BGE-M3  1024-dim"),
    ("LLM  (EN)",       "Gemini 2.5 Flash"),
    ("LLM  (AS/IN)",    "Sarvam 30B Indic"),
    ("TTS",             "Sarvam TTS  EN+AS"),
    ("Document DB",     "MongoDB Atlas"),
    ("Auth / API",      "FastAPI · JWT RS256"),
]
sx_start = 0.18
sx_bw    = 1.57
sx_gap   = 0.06
for i, (top, bot) in enumerate(stack_items):
    lx = sx_start + i * (sx_bw + sx_gap)
    add_rounded_rect(s1, lx, stack_y + 0.26, sx_bw, 0.48, RGBColor(0x0C,0x11,0x1E), BORDER, 0.7)
    add_text_box(s1, lx, stack_y + 0.28, sx_bw, 0.20,
                 top, font_size=6.5, bold=True, color=MIDGREY)
    add_text_box(s1, lx, stack_y + 0.48, sx_bw, 0.20,
                 bot, font_size=7, bold=False, color=LTGREY)


# ── Footer slide 1 ──────────────────────────────────────────────────────────
add_rect(s1, 0, 7.26, 13.33, 0.24, DARKGREY)
add_text_box(s1, 0.2, 7.28, 13, 0.20,
             "Slide 1 of 2  ·  AI-Powered Bilingual Educational Platform  ·  System Architecture  ·  MSME IdeaHackathon 6.0",
             font_size=6.5, bold=False, color=MIDGREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TECHNICAL PROCESSING BRIEF
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_rect(s2, 0, 0, 13.33, 7.5, BG)
add_rect(s2, 0, 0, 13.33, 0.52, DARKGREY)
add_text_box(s2, 0.2, 0.04, 13, 0.42,
             "BILINGUAL AI EDUCATIONAL PLATFORM  ·  TECHNICAL PROCESSING BRIEF",
             font_size=11, bold=True, color=LTGREY)

# ── LEFT COLUMN: Pipeline Detail Table ────────────────────────────────────────
add_text_box(s2, 0.20, 0.60, 6.20, 0.26,
             "PIPELINE STAGE SPECIFICATIONS", font_size=8, bold=True, color=CYAN)

table_data = [
    ("STAGE",               "MODULE / ENGINE",              "INPUT",                    "OUTPUT",                  "KEY PARAMETERS"),
    ("1. Input Reception",  "Web / Mobile API Gateway",     "HTTP req · voice blob",    "Normalised query string", "lang_hint · session_id"),
    ("2. Language ID",      "Statistical LID Classifier",   "Raw query string",         "lang_code ∈ {en, as}",    "conf_threshold = 0.85"),
    ("3. Query Embedding",  "BGE-M3  (1024-dim)",           "Query string",             "float[1024] vector",      "model=bge-m3; pooling=cls"),
    ("4. Topic Matching",   "Cosine Similarity Engine",     "Query vector + topic DB",  "topic_id, score",         "n_topics=197; threshold=0.72"),
    ("5. Scope Resolution", "Curriculum Resolver Service",  "topic_id",                 "chapter·subject·board",   "hierarchy depth=6 levels"),
    ("6. RAG Retrieval",    "HNSW Vector Index",            "Query vector + filters",   "Top-K text chunks",       "k=5; ef_search=128"),
    ("7. Context Assembly", "Prompt Builder",               "chunks + memory + scope",  "Structured prompt string","max_ctx=12 000 tokens"),
    ("8. LLM Generation",   "Gemini Flash / Sarvam 30B",    "Prompt string",            "Answer text stream",      "max_new_tokens=800; temp=0.3"),
    ("9. Response Delivery","SSE Stream Handler + TTS",     "Answer text",              "JSON stream + audio",     "format=SSE; voice=indic-f1"),
    ("10. Memory Write",    "Student Profile Store",        "Session summary",          "Updated memory doc",      "TTL=∞; upsert=true"),
]

col_widths = [1.38, 1.55, 1.40, 1.40, 1.42]
col_starts = [0.20]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w + 0.02)

row_h   = 0.51
start_ty = 0.90

for r, row in enumerate(table_data):
    row_y = start_ty + r * row_h * 0.56
    for c, cell in enumerate(row):
        lx = col_starts[c]
        cw = col_widths[c]
        if r == 0:
            fill_c = RGBColor(0x1A,0x23,0x42)
            txt_c  = CYAN
            bold_c = True
            fs     = 6.5
        elif r % 2 == 0:
            fill_c = RGBColor(0x0C,0x11,0x1E)
            txt_c  = LTGREY
            bold_c = (c == 0)
            fs     = 6.5
        else:
            fill_c = RGBColor(0x0E,0x14,0x24)
            txt_c  = LTGREY
            bold_c = (c == 0)
            fs     = 6.5

        add_rect(s2, lx, row_y, cw, row_h * 0.54, fill_c)
        add_text_box(s2, lx + 0.04, row_y + 0.01, cw - 0.06, row_h * 0.52,
                     cell, font_size=fs, bold=bold_c, color=txt_c, align=PP_ALIGN.LEFT, wrap=True)


# ── RIGHT COLUMN: Innovation Differentiators ───────────────────────────────────
rx = 7.22
add_text_box(s2, rx, 0.60, 5.95, 0.26,
             "TECHNICAL DIFFERENTIATORS", font_size=8, bold=True, color=INDIGO)

diffs = [
    (INDIGO, "01 · Dual-Language RAG",
     "Single retrieval pipeline serves both English and Assamese queries from the same indexed corpus. "
     "Language-conditioned prompting switches the generation language without maintaining separate indices. "
     "Reduces storage by ~40 % vs. dual-index approach."),
    (GREEN, "02 · Curriculum-Scoped Retrieval",
     "Metadata filters (board_id, class_id, subject_id, chapter_id) applied at vector-search time, not post-hoc. "
     "Ensures retrieved chunks are syllabus-relevant, reducing hallucinations to near-zero on in-scope topics. "
     "Filter cardinality: 6-level hierarchy × 3 boards × 4 classes."),
    (PURPLE, "03 · Persistent Student Memory",
     "MongoDB document per student accumulates topics_studied, weak_areas, doubt_history. "
     "Injected into every prompt as a 'prior context' block (max 800 tokens). "
     "Enables multi-session personalisation unavailable in stateless LLM deployments."),
    (AMBER, "04 · Human-in-the-Loop Content QC",
     "Staff publish pipeline: raw_content → chunker → translator → embedder → QC review → live. "
     "No auto-published AI-generated curriculum content. "
     "Version-controlled with optimistic locking (version counter per chapter)."),
    (CYAN, "05 · Edge-First Architecture",
     "Cloudflare Workers handle routing, JWT auth, CORS, rate-limiting, and embedding calls at the edge. "
     "P95 TTFB < 400 ms for cached responses. "
     "Origin (Cloud Run) invoked only for LLM generation and DB writes."),
]

dy = 0.90
for (col, title, body) in diffs:
    dark_map2 = {
        INDIGO: RGBColor(0x0E,0x0C,0x22), GREEN: RGBColor(0x05,0x16,0x10),
        PURPLE: RGBColor(0x10,0x06,0x28), AMBER: RGBColor(0x16,0x0E,0x00),
        CYAN:   RGBColor(0x04,0x18,0x20),
    }
    bh = 0.85
    add_rounded_rect(s2, rx, dy, 5.95, bh, dark_map2[col], col, 1.2)
    add_text_box(s2, rx + 0.10, dy + 0.04, 5.75, 0.24,
                 title, font_size=8, bold=True, color=col)
    add_text_box(s2, rx + 0.10, dy + 0.28, 5.75, bh - 0.30,
                 body, font_size=6.8, bold=False, color=LTGREY, align=PP_ALIGN.LEFT)
    dy += bh + 0.06


# ── Bottom: System-level constraints ─────────────────────────────────────────
add_text_box(s2, 0.20, 6.26, 13.0, 0.22,
             "SYSTEM CONSTRAINTS & PERFORMANCE TARGETS", font_size=7, bold=True, color=AMBER)

constraints = [
    ("Latency",          "TTFB < 400 ms (edge cache)\nLLM stream start < 2 s"),
    ("Throughput",       "1 000 concurrent users\nEdge: stateless horizontal scale"),
    ("Accuracy",         "RAG grounding: in-scope topics\nhallucination rate < 2 %"),
    ("Language support", "English (en) + Assamese (as)\nExtensible to 22 IN languages"),
    ("Content coverage", "SEBA Cl.9–10 · AHSEC Cl.11–12\nCore subjects + PYQ archive"),
    ("Auth & Security",  "JWT RS256 · rate limiting\nHTTPS-only · CORS enforced"),
    ("Data privacy",     "Student data stored regionally\nNo PII in LLM prompts"),
    ("Availability",     "Cloud Run + Cloudflare CDN\nTargeted 99.5 % monthly uptime"),
]

csx   = 0.20
cs_bw = 1.55
cs_gap = 0.06
cs_y  = 6.52
for i, (label, val) in enumerate(constraints):
    lx = csx + i * (cs_bw + cs_gap)
    add_rounded_rect(s2, lx, cs_y, cs_bw, 0.70, RGBColor(0x0A,0x10,0x1C), BORDER, 0.7)
    add_text_box(s2, lx + 0.05, cs_y + 0.03, cs_bw - 0.08, 0.22,
                 label, font_size=6.5, bold=True, color=AMBER)
    add_text_box(s2, lx + 0.05, cs_y + 0.26, cs_bw - 0.08, 0.40,
                 val, font_size=6.5, bold=False, color=LTGREY, align=PP_ALIGN.LEFT)


# ── Footer slide 2 ────────────────────────────────────────────────────────────
add_rect(s2, 0, 7.26, 13.33, 0.24, DARKGREY)
add_text_box(s2, 0.2, 7.28, 13, 0.20,
             "Slide 2 of 2  ·  AI-Powered Bilingual Educational Platform  ·  Technical Brief  ·  MSME IdeaHackathon 6.0",
             font_size=6.5, bold=False, color=MIDGREY)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/runner/workspace/attached_assets/Syrabit_Technical_Block_Diagram.pptx"
prs.save(out)
print(f"Saved: {out}")
import os
print(f"Size: {os.path.getsize(out):,} bytes ({os.path.getsize(out)/1024:.1f} KB)")
